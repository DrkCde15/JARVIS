import os
import json
import subprocess
from pathlib import Path
import pandas as pd
import fitz
from docx import Document
from pptx import Presentation
from commands.constants import Colors
from commands.voice import falar
from ai_service import gerar_resposta_ia, extrair_params_ia
from cli_design import jarvis_ask

def encontrar_pasta(nome_pasta_usuario):
    home = os.path.expanduser("~")
    mapeamento = {
        "documentos": "Documents",
        "downloads": "Downloads",
        "imagens": "Pictures",
        "desktop": "Desktop",
        "musicas": "Music",
        "videos": "Videos"
    }
    nome_sistema = mapeamento.get(nome_pasta_usuario.lower(), nome_pasta_usuario)
    caminho = os.path.join(home, nome_sistema)
    if os.path.exists(caminho):
        return caminho
    for item in os.listdir(home):
        if nome_pasta_usuario.lower() in item.lower():
            caminho_tentativa = os.path.join(home, item)
            if os.path.isdir(caminho_tentativa):
                return caminho_tentativa
    return None

def abrir_pasta(match, username):
    nome_pasta_usuario = match.group(1).strip() if hasattr(match, 'group') else match
    caminho = encontrar_pasta(nome_pasta_usuario)
    if caminho:
        try:
            comando = f"Start-Process explorer -ArgumentList '{caminho}' -WindowStyle Maximized"
            subprocess.run(["powershell", "-Command", comando], check=True)
            return f"Abrindo a pasta {os.path.basename(caminho)}, senhor."
        except Exception as e:
            return f"Erro ao abrir a pasta: {e}"
    return f"Pasta '{nome_pasta_usuario}' não encontrada, senhor."

def criar_arquivo(match, username, status=None):
    documentos = Path.home() / "Documents"
    texto_original = match if isinstance(match, str) else (match.group(0) if hasattr(match, 'group') else "")
    params = extrair_params_ia(texto_original, ["nome", "conteudo"], username=username) if texto_original else {}

    nome = params.get("nome") or ""
    conteudo = params.get("conteudo") or ""

    if not nome:
        nome = jarvis_ask("Como devo chamar este arquivo de texto, senhor? Por exemplo: notas.txt", status)
    if not nome:
        return "Operação cancelada — nenhum nome fornecido."
    if not conteudo:
        conteudo = jarvis_ask("Certo. Qual será o conteúdo dele?", status)

    caminho_completo = documentos / nome
    try:
        with open(caminho_completo, "w", encoding="utf-8") as f:
            f.write(conteudo)
        return f"✅ Arquivo '{nome}' criado com sucesso na pasta Documentos."
    except Exception as e:
        return f"❌ Erro ao criar o arquivo: {e}"

def criar_codigo(match, username, session_id=None, status=None):
    documentos = Path.home() / "Documents"
    documentos.mkdir(exist_ok=True)

    texto_original = match if isinstance(match, str) else (match.group(0) if hasattr(match, 'group') else "")
    params = extrair_params_ia(
        texto_original,
        ["linguagem", "descricao", "nome_arquivo"],
        username=username,
    ) if texto_original else {}

    linguagem = (params.get("linguagem") or "").lower()
    descricao = params.get("descricao") or ""
    nome_base = params.get("nome_arquivo") or ""

    if not linguagem:
        linguagem = jarvis_ask("Senhor, qual linguagem de programação você quer usar? Ex: Python, JavaScript, Go...", status).lower()
    if not linguagem:
        return "Operação cancelada — nenhuma linguagem informada."

    if not descricao:
        descricao = jarvis_ask("Perfeito. Descreva o que esse código deve fazer, por favor.", status)
    if not descricao:
        return "Operação cancelada — nenhuma descrição fornecida."

    prompt = f"Crie um código em {linguagem} que: {descricao}"
    try:
        codigo = gerar_resposta_ia(prompt, session_id, username or "Senhor")
    except Exception as e:
        return f"❌ Erro ao gerar código com IA: {e}"

    extensoes = {
        "python": ".py", "javascript": ".js", "java": ".java", "html": ".html",
        "c": ".c", "c++": ".cpp", "go": ".go", "php": ".php", "ruby": ".rb",
        "kotlin": ".kt", "swift": ".swift", "rust": ".rs", "csharp": ".cs",
        "c#": ".cs", "css": ".css", "sql": ".sql", "r": ".r"
    }

    ext = extensoes.get(linguagem, ".txt")
    if not nome_base:
        nome_base = jarvis_ask("Qual nome devo dar a este arquivo? (sem extensão)", status) or "codigo_gerado"
    nome_arquivo = nome_base.strip() + ext
    caminho = documentos / nome_arquivo
    try:
        with open(caminho, "w", encoding="utf-8") as f:
            f.write(codigo)
        return f"✅ Código gerado e salvo em: {caminho}"
    except Exception as e:
        return f"❌ Erro ao salvar o arquivo: {e}"

def listar_arquivos(match, username):
    extensao = (match.group(1) or "").lower().replace(".", "").strip() if hasattr(match, 'group') else ""
    pasta = (match.group(2) or "Documentos").lower().strip() if hasattr(match, 'group') and match.lastindex >= 2 else "Documentos"
    
    base_path = Path.home()
    nome_pasta = {
        "documentos": "Documents",
        "area de trabalho": "Desktop",
        "área de trabalho": "Desktop",
        "desktop": "Desktop",
        "downloads": "Downloads"
    }.get(pasta, pasta)
    
    diretorio = base_path / nome_pasta
    if not diretorio.exists():
        return f"A pasta '{nome_pasta}' não foi encontrada, senhor."
    
    try:
        arquivos = list(diretorio.rglob(f"*.{extensao}")) if extensao else list(diretorio.rglob("*"))
        if not arquivos:
            return f"Senhor, a pasta '{nome_pasta}' está vazia ou não contém arquivos '.{extensao}'."
        
        lista = "\n- " + "\n- ".join([str(arq.relative_to(base_path)) for arq in arquivos[:15]])
        return f"Encontrei arquivos em '{nome_pasta}':\n{lista}"
    except Exception as e:
        return f"Erro ao listar arquivos: {e}"

# ========== Leitura de arquivos ==========

CODE_EXTENSIONS = {
    ".py", ".js", ".ts", ".tsx", ".jsx", ".mjs", ".cjs",
    ".java", ".c", ".cpp", ".cc", ".cxx", ".h", ".hpp", ".hxx",
    ".go", ".rs", ".rb", ".php", ".swift", ".kt", ".kts", ".scala",
    ".cs", ".fs", ".vb", ".r", ".m", ".mm",
    ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf",
    ".html", ".css", ".scss", ".less", ".sass",
    ".sql", ".sh", ".ps1", ".bat", ".cmd",
    ".xml", ".svg", ".json", ".md",
    ".dockerfile", ".tf", ".gradle", ".cmake",
}

LINGUAGENS_POR_EXTENSAO = {
    ".py": "Python", ".js": "JavaScript", ".ts": "TypeScript",
    ".tsx": "TypeScript React", ".jsx": "JavaScript React",
    ".java": "Java", ".c": "C", ".cpp": "C++", ".h": "C/C++ Header",
    ".go": "Go", ".rs": "Rust", ".rb": "Ruby", ".php": "PHP",
    ".swift": "Swift", ".kt": "Kotlin", ".cs": "C#",
    ".yaml": "YAML", ".yml": "YAML", ".toml": "TOML",
    ".html": "HTML", ".css": "CSS", ".scss": "SCSS",
    ".sql": "SQL", ".sh": "Shell Script", ".ps1": "PowerShell",
    ".json": "JSON", ".md": "Markdown", ".xml": "XML",
    ".dockerfile": "Dockerfile", ".tf": "Terraform",
}


def _confirmar_analise(match, username):
    nome = match.group(1).strip() if hasattr(match, 'group') else match
    resposta = jarvis_ask(
        f"Deseja que eu analise '{nome}'? "
        "Posso ler o arquivo e, se for código, executar ferramentas como linters e testes. "
        "Digite SIM para autorizar.",
    )
    return resposta.strip().lower() in {"sim", "s", "yes", "y"}


def _analisar_com_powershell(arquivo: Path, session_id: str, username: str):
    """Executa ferramentas de análise (lint, typecheck, teste) via PowerShell."""
    sufixo = arquivo.suffix.lower()
    resultados = []

    if sufixo == ".py":
        resposta = jarvis_ask(
            "Deseja rodar linters e testes neste arquivo Python? "
            "Posso executar ruff, mypy e pytest. Digite SIM para autorizar ou NAO para pular.",
        )
        if resposta.strip().lower() in {"sim", "s", "yes", "y"}:
            for cmd, label in [
                (f'python -m ruff check "{arquivo}" 2>&1', "ruff (lint)"),
                (f'python -m mypy "{arquivo}" 2>&1', "mypy (tipos)"),
                (f'python -m pytest "{arquivo}" -v --tb=short 2>&1', "pytest (testes)"),
            ]:
                try:
                    r = subprocess.run(
                        ["powershell", "-NoProfile", "-Command", cmd],
                        capture_output=True, text=True, timeout=30, check=False,
                    )
                    saida = (r.stdout.strip() + "\n" + r.stderr.strip()).strip()
                    if saida:
                        resultados.append(f"--- {label} ---\n{saida[:1500]}")
                except subprocess.TimeoutExpired:
                    resultados.append(f"--- {label} ---\n(tempo excedido)")
                except Exception as e:
                    resultados.append(f"--- {label} ---\n(erro: {e})")
    elif sufixo in (".js", ".ts", ".tsx", ".jsx"):
        resposta = jarvis_ask(
            "Deseja rodar ESLint neste arquivo? Digite SIM para autorizar ou NAO para pular.",
        )
        if resposta.strip().lower() in {"sim", "s", "yes", "y"}:
            try:
                result = subprocess.run(
                    ["powershell", "-NoProfile", "-Command",
                     f'npx eslint "{arquivo}" 2>&1'],
                    capture_output=True, text=True, timeout=30, check=False,
                )
                saida = (result.stdout.strip() + "\n" + result.stderr.strip()).strip()
                if saida:
                    resultados.append(f"--- ESLint ---\n{saida[:1500]}")
            except Exception:
                pass

    return resultados


def ler_codigo(caminho: Path) -> str | None:
    try:
        with open(caminho, "r", encoding="utf-8") as f:
            return f.read()
    except Exception:
        try:
            with open(caminho, "r", encoding="latin-1") as f:
                return f.read()
        except Exception:
            return None

def buscar_arquivo_por_nome(nome_arquivo, pasta_base=None):
    if pasta_base is None:
        pasta_base = Path.home() / "Documents"
    for arquivo in pasta_base.rglob("*"):
        if arquivo.name.lower() == nome_arquivo.lower():
            return arquivo.resolve()
    return None

def ler_docx(caminho):
    try:
        doc = Document(caminho)
        return '\n'.join([p.text for p in doc.paragraphs])
    except Exception as e: return f"Erro: {e}"

def ler_txt(caminho):
    try:
        with open(caminho, "r", encoding="utf-8") as f: return f.read()
    except Exception as e: return f"Erro: {e}"

def ler_csv(caminho):
    try:
        df = pd.read_csv(caminho)
        return df.to_string(index=False)
    except Exception as e: return f"Erro: {e}"

def ler_json(caminho):
    try:
        with open(caminho, "r", encoding="utf-8") as f: return json.dumps(json.load(f), indent=4)
    except Exception as e: return f"Erro: {e}"

def ler_pdf(caminho):
    try:
        doc = fitz.open(caminho)
        return ''.join(pagina.get_text() for pagina in doc)
    except Exception as e: return f"Erro: {e}"

def ler_excel(caminho):
    try:
        df = pd.read_excel(caminho)
        return df.to_string(index=False)
    except Exception as e: return f"Erro: {e}"

def ler_pptx(caminho):
    try:
        prs = Presentation(caminho)
        texto = ''
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"): texto += shape.text + "\n"
        return texto
    except Exception as e: return f"Erro: {e}"

def analisar_arquivos(match, username, session_id=None, permitir_powershell=True):
    try:
        nome_arquivo = match.group(1).strip() if hasattr(match, 'group') else match
        arquivo_encontrado = buscar_arquivo_por_nome(nome_arquivo)
        if not arquivo_encontrado:
            return f"Não encontrei o arquivo '{nome_arquivo}' na pasta Documentos, senhor."

        if not _confirmar_analise(match, username):
            return "Análise cancelada."

        sufixo = arquivo_encontrado.suffix.lower()
        conteudo = None
        tipo = "documento"

        if sufixo == ".txt": conteudo = ler_txt(arquivo_encontrado)
        elif sufixo == ".docx": conteudo = ler_docx(arquivo_encontrado)
        elif sufixo == ".csv": conteudo = ler_csv(arquivo_encontrado)
        elif sufixo == ".json": conteudo = ler_json(arquivo_encontrado)
        elif sufixo == ".pdf": conteudo = ler_pdf(arquivo_encontrado)
        elif sufixo in [".xlsx", ".xls"]: conteudo = ler_excel(arquivo_encontrado)
        elif sufixo == ".pptx": conteudo = ler_pptx(arquivo_encontrado)
        elif sufixo in CODE_EXTENSIONS:
            conteudo = ler_codigo(arquivo_encontrado)
            tipo = LINGUAGENS_POR_EXTENSAO.get(sufixo, "código")
        else:
            return f"Formato '{sufixo}' não suportado."

        if not conteudo or not conteudo.strip() or "Erro:" in str(conteudo):
            return "O arquivo está vazio ou ilegível, senhor."

        # Análise PowerShell para código
        resultados_ps = []
        if tipo != "documento" and permitir_powershell:
            resultados_ps = _analisar_com_powershell(arquivo_encontrado, session_id, username)

        extras = ""
        if resultados_ps:
            extras = "\n\n--- Resultados de ferramentas ---\n" + "\n\n".join(resultados_ps)

        linhas = conteudo.count("\n") + 1
        prompt = (
            f"Analise este arquivo de {tipo}:\n"
            f"Caminho: {arquivo_encontrado}\n"
            f"Linhas: {linhas}\n\n"
            f"```\n{conteudo[:8000]}\n```"
            f"{extras}"
        )
        if len(conteudo) > 8000:
            prompt += (
                f"\n\n[O arquivo tem {linhas} linhas. Mostrei as primeiras ~8000 chars. "
                f"Se precisar de mais contexto, peça que leio partes específicas.]"
            )

        return gerar_resposta_ia(prompt, session_id, username or "Senhor")
    except Exception as e:
        return f"Erro ao analisar arquivo: {e}"
