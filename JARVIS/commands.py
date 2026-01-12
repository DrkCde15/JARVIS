import json
import os
import re
import unicodedata
import subprocess
import requests
import sys
import ctypes
import shutil
from pathlib import Path
import webbrowser
import pyautogui # type: ignore
from bs4 import BeautifulSoup
import yt_dlp # type: ignore
import time
from datetime import datetime
import pandas as pd
from PIL import Image
from ai_service import (
    MODEL_NAME,
    recarregar_llm,
    construir_historico,
    gerar_resposta_ia
)
from memory import (
    adicionar_mensagem_chat,
    registrar_log,
    salvar_senha_smtp,
    obter_senha_smtp,
    obter_session_id_por_token
)
import fitz
from docx import Document
from pptx import Presentation
import pyttsx3
import threading
from queue import Queue
from typing import Callable
import warnings
import pywhatkit as kit # type: ignore
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from email.mime.base import MIMEBase
from email import encoders
from getpass import getpass
import winapps # type: ignore
from plyer import notification  # type: ignore

warnings.filterwarnings('ignore')
if "USER_AGENT" not in os.environ:
    os.environ["USER_AGENT"] = "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/115.0.0.0 Safari/537.36"

# ========== Cores ==========
class Colors:
    """Códigos ANSI para cores e estilos"""
    BLUE = '\033[38;5;39m'
    CYAN = '\033[38;5;51m'
    PURPLE = '\033[38;5;141m'
    MAGENTA = '\033[38;5;199m'
    PINK = '\033[38;5;213m'
    GRAY = '\033[38;5;240m'
    WHITE = '\033[97m'
    GREEN = '\033[92m'
    YELLOW = '\033[93m'
    RED = '\033[91m'
    ORANGE = '\033[38;5;208m'
    BOLD = '\033[1m'
    DIM = '\033[2m'
    RESET = '\033[0m'
    CLEAR_LINE = '\033[2K'

# ========== Configurações de voz ==========
class VoiceCommandSystem:
    def __init__(self):
        self.engine = self._init_voice_engine()
        self.command_queue = Queue()
        self.voice_lock = threading.Lock()
        self.is_speaking = False
        self._start_command_processor()

    def _init_voice_engine(self):
        try:
            engine = pyttsx3.init()
            engine.setProperty('rate', 180)
            engine.setProperty('volume', 0.9)
            return engine
        except Exception as e:
            print(f"Voice engine initialization failed: {e}")
            return None

    def _start_command_processor(self):
        def processor():
            while True:
                cmd_func, args, kwargs = self.command_queue.get()
                try:
                    time.sleep(0.3)
                    cmd_func(*args, **kwargs)
                except Exception as e:
                    print(f"Command execution error: {e}")
                finally:
                    self.command_queue.task_done()

        threading.Thread(target=processor, daemon=True).start()

    def speak(self, text: str):
        """Non-blocking speech with queue management"""
        print(f"JARVIS: {text}")
        
        if not self.engine:
            return

        def _speak():
            with self.voice_lock:
                self.is_speaking = True
                try:
                    self.engine.say(text)
                    self.engine.runAndWait()
                except Exception as e:
                    print(f"Speech error: {e}")
                finally:
                    self.is_speaking = False

        threading.Thread(target=_speak, daemon=True).start()

    def add_command(self, command_func: Callable, *args, **kwargs):
        """Add command to execution queue"""
        self.command_queue.put((command_func, args, kwargs))

# Global voice system instance
voice_system = VoiceCommandSystem()

def falar(texto: str):
    """Improved speak function that uses the enhanced voice system"""
    voice_system.speak(texto)

#=========== Funções de Permissão ==========

def is_admin():
    """Checa se o script está rodando como administrador"""
    try:
        return ctypes.windll.shell32.IsUserAnAdmin()
    except:
        return False

def relancar_como_admin():
    """Relança o script atual como administrador"""
    if not is_admin():
        print("Tentando reiniciar como admin...")
        ret = ctypes.windll.shell32.ShellExecuteW(
            None, "runas", sys.executable, ' '.join(sys.argv), None, 1
        )
        if ret <= 32:
            print("Falha ao tentar executar como administrador.")
            return False
        print("Script relançado como administrador, finalizando instância atual.")
        sys.exit(0)
    return True

#=========== Chocolatey ==========

def verificar_choco_instalado():
    """Verifica se o Chocolatey está instalado"""
    try:
        resultado = subprocess.run(["choco", "-v"], capture_output=True, text=True)
        if "Chocolatey" in resultado.stdout or resultado.returncode == 0:
            return True
    except FileNotFoundError:
        return False
    return False

def instalar_chocolatey_via_powershell():
    """Instala o Chocolatey via PowerShell, sem pyautogui"""
    print("[*] Instalando Chocolatey via PowerShell...")
    comando = (
        'Set-ExecutionPolicy Bypass -Scope Process -Force; '
        '[System.Net.ServicePointManager]::SecurityProtocol = '
        '[System.Net.ServicePointManager]::SecurityProtocol -bor 3072; '
        'iex ((New-Object System.Net.WebClient).DownloadString("https://chocolatey.org/install.ps1"))'
    )
    processo = subprocess.run(["powershell", "-Command", comando], capture_output=True, text=True)
    if processo.returncode == 0:
        print("[+] Chocolatey instalado com sucesso.")
        return True
    else:
        print("[-] Falha ao instalar Chocolatey.")
        print(processo.stderr)
        return False

def instalar_programa_choco(programa):
    """Instala o programa via Chocolatey no terminal, usando subprocess"""
    comando = f"choco install {programa} -y"
    print(f"[*] Instalando {programa} via Chocolatey...")
    processo = subprocess.run(comando, shell=True)
    if processo.returncode == 0:
        print(f"[+] {programa} instalado com sucesso.")
        return True
    else:
        print(f"[-] Falha ao instalar {programa}.")
        return False

def instalar_programa_via_cmd_admin(programa=None, username=None):
    """Função principal para instalar programas via choco com privilégios admin"""
    if not programa:
        return "Nenhum programa informado para instalação."

    if not relancar_como_admin():
        return "Erro ao tentar elevar privilégios."

    if not verificar_choco_instalado():
        print("[*] Chocolatey não encontrado, instalando...")
        sucesso = instalar_chocolatey_via_powershell()
        if not sucesso:
            return "Falha ao instalar Chocolatey, abortando."

        time.sleep(5)

    sucesso = instalar_programa_choco(programa)
    if sucesso:
        return f"Programa {programa} instalado com sucesso."
    else:
        return f"Falha ao instalar o programa {programa}."

#=========== Desinstalação de Programa ==========
def desinstalar_programa(nome_programa, username, modo='texto'):
    try:
        resposta = f"[*] Desinstalando {nome_programa} via Chocolatey..."
        if modo == 'voz':
            falar(resposta)
        print(resposta)

        comando = f'choco uninstall {nome_programa} -y'
        subprocess.run(comando, shell=True)
        return f"{nome_programa} desinstalado com sucesso (ou já não estava instalado)."
    except Exception as e:
        return f"Erro ao desinstalar {nome_programa}: {e}"

# ========== Gerenciamento de Aplicativos ==========

def listar_aplicativos_winapps(match=None, username=None):
    """Lista todos os aplicativos instalados usando winapps"""
    try:
        apps = list(winapps.list_installed())
        if not apps:
            return "Nenhum aplicativo encontrado, senhor."
        
        lista = ["Aplicativos instalados:"]
        for app in apps[:100]:  # Limita a 100
            versao = app.version if app.version else 'N/A'
            lista.append(f"- {app.name} (v{versao})")

        if len(apps) > 100:
            lista.append(f"\n... e mais {len(apps) - 100} aplicativos.")
        
        return "\n".join(lista)
    except Exception as e:
        return f"Erro ao listar aplicativos: {e}"


def buscar_aplicativo_winapps(nome_app):
    """Busca um aplicativo específico instalado"""
    try:
        for app in winapps.search_installed(nome_app):
            return app
        return None
    except Exception as e:
        print(f"Erro ao buscar aplicativo: {e}")
        return None

def abrir_aplicativo_winapps(match, username=None):
    """Abre aplicativo usando winapps"""
    nome = match.group(2).lower().strip()
    
    try:
        app = buscar_aplicativo_winapps(nome)
        
        if app:
            # Tentar abrir via local de instalação
            if app.install_location and os.path.exists(app.install_location):
                subprocess.Popen(f'start "" "{app.install_location}"', shell=True)
            # Se não tiver local, tentar abrir pelo nome
            else:
                subprocess.Popen(f'start {app.name}', shell=True)
            return f"Abrindo {app.name}, senhor."
        else:
            return f"Aplicativo '{nome}' não encontrado, senhor."
            
    except Exception as e:
        return f"Erro ao abrir aplicativo: {e}"


def desinstalar_app_winapps(nome_app, username=None, modo='texto'):
    """Desinstala aplicativo usando winapps"""
    try:
        app = buscar_aplicativo_winapps(nome_app)
        
        if not app:
            return f"Aplicativo '{nome_app}' não encontrado, senhor."
        
        resposta = f"Desinstalando {app.name}..."
        if modo == 'voz':
            falar(resposta)
        print(resposta)
        
        # Usar chocolatey como fallback se winapps não conseguir desinstalar
        try:
            app.uninstall(wait=True)
            return f"{app.name} desinstalado com sucesso, senhor."
        except:
            return desinstalar_programa(nome_app, username, modo)
        
    except Exception as e:
        # Fallback para chocolatey
        return desinstalar_programa(nome_app, username, modo)


def info_aplicativo_winapps(nome_app, username=None):
    """Obtém informações detalhadas de um aplicativo"""
    try:
        app = buscar_aplicativo_winapps(nome_app)
        
        if not app:
            return f"Aplicativo '{nome_app}' não encontrado, senhor."
        
        info = f"""Informações do Aplicativo:
- Nome: {app.name}
- Versão: {app.version if app.version else 'N/A'}
- Editor: {app.publisher if app.publisher else 'N/A'}
- Data de Instalação: {app.install_date if app.install_date else 'N/A'}
- Local: {app.install_location if app.install_location else 'N/A'}"""
        
        return info.strip()
        
    except Exception as e:
        return f"Erro ao obter informações: {e}"


# ========== Automação de WhatsApp ==========

def enviar_whatsapp_agendado(match, username=None, modo='texto'):
    """Envia mensagem no WhatsApp com horário agendado"""
    try:
        numero = input(f"{Colors.PURPLE}>{Colors.RESET} Digite o número com DDI (ex: +5511999999999): ").strip()
        mensagem = input(f"{Colors.PURPLE}>{Colors.RESET} Digite a mensagem: ").strip()
        
        # Agenda para pelo menos 2 minutos à frente
        agora = datetime.now()
        
        # Garantir que temos pelo menos 2 minutos no futuro
        minuto_envio = agora.minute + 2
        hora_envio = agora.hour
        
        # Ajustar se passar de 59 minutos
        if minuto_envio >= 60:
            minuto_envio = minuto_envio - 60
            hora_envio = hora_envio + 1
            if hora_envio >= 24:
                hora_envio = 0
        
        msg = f"Enviando mensagem para {numero} às {hora_envio:02d}:{minuto_envio:02d}..."
        if modo == 'voz':
            falar(msg)
        print(msg)
        
        # Usar wait_time mínimo de 15 segundos
        kit.sendwhatmsg(numero, mensagem, hora_envio, minuto_envio, wait_time=15, tab_close=True)
        
        return f"Mensagem agendada com sucesso para {numero} às {hora_envio:02d}:{minuto_envio:02d}, senhor."
        
    except Exception as e:
        erro_msg = f"Erro ao enviar WhatsApp: {e}"
        if "sleep length must be non-negative" in str(e):
            # Tentar com mais tempo no futuro
            try:
                agora = datetime.now()
                minuto_envio = agora.minute + 5  # 5 minutos no futuro
                hora_envio = agora.hour
                
                if minuto_envio >= 60:
                    minuto_envio = minuto_envio - 60
                    hora_envio = hora_envio + 1
                
                kit.sendwhatmsg(numero, mensagem, hora_envio, minuto_envio, wait_time=15, tab_close=False)
                return f"Mensagem reagendada para {hora_envio:02d}:{minuto_envio:02d}, senhor."
            except:
                return "Erro: Não foi possível agendar a mensagem. Tente novamente mais tarde."
        return erro_msg


def enviar_whatsapp(match, username=None, modo='texto'):
    """Envia mensagem instantânea no WhatsApp"""
    try:
        numero = input(f"{Colors.PURPLE}>{Colors.RESET} Digite o número com DDI (ex: +5511999999999): ").strip()
        mensagem = input(f"{Colors.PURPLE}>{Colors.RESET} Digite a mensagem: ").strip()
        
        msg = f"Enviando mensagem instantânea para {numero}..."
        if modo == 'voz':
            falar(msg)
        print(msg)
        
        # Método 1: Tentar sendwhatmsg_instantly com wait_time adequado
        try:
            # wait_time mínimo de 10 segundos para evitar erro negativo
            kit.sendwhatmsg_instantly(numero, mensagem, wait_time=10, tab_close=True)
            return f"Mensagem enviada instantaneamente para {numero}, senhor."
        except Exception as e1:
            if "sleep length must be non-negative" in str(e1):
                # Método 2: Usar sendwhatmsg com 1 minuto no futuro
                agora = datetime.now()
                minuto_envio = agora.minute + 1
                hora_envio = agora.hour
                
                if minuto_envio >= 60:
                    minuto_envio = 0
                    hora_envio = hora_envio + 1
                    if hora_envio >= 24:
                        hora_envio = 0
                
                kit.sendwhatmsg(numero, mensagem, hora_envio, minuto_envio, wait_time=10, tab_close=False)
                return f"Mensagem agendada para {hora_envio:02d}:{minuto_envio:02d} (quase instantâneo), senhor."
            else:
                raise e1
        
    except Exception as e:
        erro_msg = f"Erro ao enviar WhatsApp: {e}"
        return erro_msg


def enviar_whatsapp_grupo(match, username=None, modo='texto'):
    """Envia mensagem para grupo usando ID obtido via Inspecionar Elemento"""
    try:
        print(f"\n{Colors.YELLOW}🔍 {Colors.BOLD}Como obter o ID do grupo:{Colors.RESET}")
        print(f"{Colors.GRAY}1.{Colors.RESET} Abra o grupo no WhatsApp Web")
        print(f"{Colors.GRAY}2.{Colors.RESET} Clique com o botão direito no nome do grupo")
        print(f"{Colors.GRAY}3.{Colors.RESET} Selecione 'Inspecionar' ou 'Inspect'")
        print(f"{Colors.GRAY}4.{Colors.RESET} Procure no código HTML por:")
        print(f"   {Colors.CYAN}data-id{Colors.RESET} ou {Colors.CYAN}data-group-id{Colors.RESET}")
        print(f"{Colors.GRAY}5.{Colors.RESET} Copie o valor, exemplo:")
        print(f"   {Colors.GREEN}120363012345678901@g.us{Colors.RESET}")
        print(f"\n{Colors.YELLOW}📝 {Colors.BOLD}Formato do ID:{Colors.RESET}")
        print(f"• {Colors.CYAN}120363012345678901@g.us{Colors.RESET}")
        print(f"• {Colors.CYAN}5511999999999-1623456789@g.us{Colors.RESET}")
        print(f"• Sempre termina com {Colors.YELLOW}@g.us{Colors.RESET}")
        
        print(f"\n{Colors.GRAY}{'='*60}{Colors.RESET}")
        
        # Pedir o ID do grupo
        grupo_id = input(f"\n{Colors.PURPLE}>{Colors.RESET} Cole o ID do grupo (@g.us): ").strip()
        mensagem = input(f"{Colors.PURPLE}>{Colors.RESET} Digite a mensagem: ").strip()
        
        if not grupo_id:
            return "ID do grupo é obrigatório, senhor."
        
        if not mensagem:
            return "Mensagem é obrigatória, senhor."
        
        # Verificar formato do ID
        if not grupo_id.endswith('@g.us'):
            print(f"{Colors.YELLOW}⚠{Colors.RESET}  O ID deve terminar com '@g.us'")
            confirmar = input(f"{Colors.YELLOW}❓{Colors.RESET}  Continuar mesmo assim? (s/n): ").strip().lower()
            if confirmar not in ['s', 'sim', 'y', 'yes']:
                return "Operação cancelada."
        
        msg = f"Preparando para enviar para o grupo..."
        if modo == 'voz':
            falar("Preparando mensagem para o grupo")
        print(f"\n{Colors.GREEN}✓{Colors.RESET} {msg}")
        print(f"{Colors.GRAY}ID: {grupo_id}{Colors.RESET}")
        
        # Agenda para 1 minuto no futuro
        agora = datetime.now()
        hora = agora.hour
        minuto = agora.minute + 1
        
        # Ajustar se passar de 59 minutos
        if minuto >= 60:
            minuto -= 60
            hora = (hora + 1) % 24
        
        print(f"{Colors.GRAY}⏰{Colors.RESET} Agendando para {hora:02d}:{minuto:02d}...")
        
        try:
            # Testar se o pywhatkit aceita o ID de grupo
            kit.sendwhatmsg_to_group(
                group_id=grupo_id,
                message=mensagem,
                time_hour=hora,
                time_min=minuto,
                wait_time=20,
                tab_close=False
            )
            
            return f"Mensagem agendada para o grupo às {hora:02d}:{minuto:02d}, senhor."
            
        except Exception as e:
            erro_msg = str(e)
            
            # Se for erro de grupo, tentar método alternativo
            if "group" in erro_msg.lower() or "inválido" in erro_msg.lower():
                print(f"{Colors.YELLOW}⚠{Colors.RESET}  Método de grupo falhou. Tentando método alternativo...")
                
                # Método alternativo: usar o ID como número de telefone (removendo @g.us)
                if '@g.us' in grupo_id:
                    numero_alternativo = grupo_id.replace('@g.us', '').replace('-', '')
                    
                    # Se o ID começar com 1203630 (comum para grupos), não funciona como número
                    if numero_alternativo.startswith('1203630'):
                        return f"ID de grupo não compatível. Formato detectado: {grupo_id}"
                    
                    print(f"{Colors.GRAY}Tentando com número: {numero_alternativo}{Colors.RESET}")
                    
                    # Tentar enviar como se fosse um número normal
                    kit.sendwhatmsg(
                        phone_no=numero_alternativo,
                        message=f"[PARA O GRUPO] {mensagem}",
                        time_hour=hora,
                        time_min=minuto,
                        wait_time=20,
                        tab_close=False
                    )
                    
                    return f"Mensagem enviada via método alternativo às {hora:02d}:{minuto:02d}"
            
            return f"Erro ao enviar: {erro_msg}"
        
    except Exception as e:
        return f"Erro geral: {e}"

# ========== YouTube e Pesquisa ==========

def tocar_musica_pywhatkit(match, username=None, modo='texto'):
    """Toca música no YouTube usando pywhatkit"""
    try:
        musica = input(f"{Colors.PURPLE}>{Colors.RESET} Qual música deseja ouvir, senhor? ").strip()
        
        if not musica:
            return "Nenhuma música informada, senhor."
        
        msg = f"Abrindo '{musica}' no YouTube..."
        if modo == 'voz':
            falar(msg)
        print(msg)
        
        kit.playonyt(musica)
        
        return f"Reproduzindo '{musica}' no YouTube, senhor."
        
    except Exception as e:
        return f"Erro ao abrir YouTube: {e}"


def pesquisar_google_pywhatkit(match, username=None, modo='texto'):
    """Pesquisa no Google usando pywhatkit"""
    try:
        # Tenta extrair o termo de pesquisa de diferentes grupos
        termo = None
        
        # Verifica diferentes padrões de grupos
        if match.lastindex >= 2:
            termo = match.group(2).strip()
        elif match.lastindex >= 1:
            termo = match.group(1).strip()
        else:
            # Se não conseguir extrair do match, pergunta ao usuário
            termo = input(f"{Colors.PURPLE}>{Colors.RESET} O que deseja pesquisar no Google, senhor? ").strip()
        
        if not termo:
            return "Por favor especifique o que deseja pesquisar, senhor."
        
        msg = f"Pesquisando '{termo}' no Google..."
        if modo == 'voz':
            falar(msg)
        print(msg)
        
        kit.search(termo)
        
        return f"Mostrando resultados para '{termo}', senhor."
        
    except Exception as e:
        return f"Erro ao pesquisar: {e}"


# ========== Funções de E-mail ==========
def enviar_email(match=None, username=None, modo="texto"):
    servidor = "smtp.gmail.com"
    porta = 587

    email_salvo, senha_salva = obter_senha_smtp(username)

    if email_salvo and senha_salva:
        remetente = email_salvo
        senha = senha_salva
        print("✓ Credenciais SMTP carregadas da memória")
    else:
        remetente = input("Seu e-mail Gmail: ").strip()
        print("Use SENHA DE APLICATIVO (16 caracteres)")
        senha = getpass("Senha: ")
        salvar_senha_smtp(username, remetente, senha)

    destinatario = input("Para: ").strip()
    assunto = input("Assunto: ").strip()

    print("Mensagem (linha vazia encerra):")
    linhas = []
    while True:
        linha = input("> ")
        if not linha.strip():
            break
        linhas.append(linha)

    mensagem = "\n".join(linhas) or "[Sem mensagem]"

    anexo = None
    if input("Anexar arquivo? (s/n): ").lower() in ("s", "sim"):
        caminho = input("Caminho: ").strip()
        if os.path.isfile(caminho):
            anexo = caminho

    msg = MIMEMultipart()
    msg["From"] = remetente
    msg["To"] = destinatario
    msg["Subject"] = assunto
    msg.attach(MIMEText(mensagem, "plain", "utf-8"))

    if anexo:
        with open(anexo, "rb") as f:
            part = MIMEBase("application", "octet-stream")
            part.set_payload(f.read())
            encoders.encode_base64(part)
            part.add_header(
                "Content-Disposition",
                f'attachment; filename="{os.path.basename(anexo)}"'
            )
            msg.attach(part)

    try:
        with smtplib.SMTP(servidor, porta, timeout=15) as server:
            server.starttls()
            server.login(remetente, senha)
            server.send_message(msg)

        registrar_log(username, f"E-mail enviado para {destinatario}")
        return f"✅ E-mail enviado para {destinatario}"

    except smtplib.SMTPAuthenticationError:
        return "❌ Falha de autenticação SMTP (senha do app inválida)"

    except Exception as e:
        return f"❌ Erro ao enviar e-mail: {e}"

# ========== Info Sites ==========
def raspar_site(url):
    try:
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
                          "AppleWebKit/537.36 (KHTML, like Gecko) "
                          "Chrome/115.0.0.0 Safari/537.36",
            "Accept-Language": "pt-BR,pt;q=0.9,en-US;q=0.8,en;q=0.7",
            "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,image/webp,*/*;q=0.8",
            "Referer": "https://www.google.com/"
        }

        response = requests.get(url, headers=headers, timeout=10)
        response.raise_for_status()
        html = response.text

        soup = BeautifulSoup(html, "html.parser")

        for tag in soup(["script", "style", "header", "footer", "nav", "form", "aside"]):
            tag.decompose()

        texto = soup.get_text(separator="\n")
        linhas = [linha.strip() for linha in texto.splitlines()]
        texto_limpo = "\n".join([linha for linha in linhas if linha])

        texto_final = texto_limpo[:1500] + ("\n..." if len(texto_limpo) > 1500 else "")

        return texto_final

    except Exception as e:
        return f"Erro ao raspar site: {e}"

def analisar_site(url, username=None):
    texto_limpo = raspar_site(url)
    if texto_limpo.startswith("Erro"):
        return texto_limpo

    prompt = (
        f"Analise e resuma o conteúdo do site abaixo:\n\n{texto_limpo}\n\n"
        "Forneça um resumo objetivo destacando pontos importantes."
    )

    resposta = gerar_resposta_ia([prompt], username)
    return resposta.content

# ========== Baixar Video ==========
def converter_audio_para_aac(caminho_video: Path):
    caminho_saida = caminho_video.with_name(caminho_video.stem + '_aac.mp4')
    cmd = [
        'ffmpeg',
        '-i', str(caminho_video),
        '-c:v', 'copy',
        '-c:a', 'aac',
        '-y',
        str(caminho_saida)
    ]
    processo = subprocess.run(cmd, capture_output=True, text=True)
    if processo.returncode != 0:
        raise RuntimeError(f'Erro ffmpeg: {processo.stderr}')
    return caminho_saida

def baixar_video_youtube(url, username, modo='texto'):
    try:
        destino = Path.home() / "Documents" / "Vídeos Download"
        destino.mkdir(parents=True, exist_ok=True)

        opcoes = {
            'format': 'bestvideo+bestaudio/best',
            'outtmpl': str(destino / '%(title)s.%(ext)s'),
            'merge_output_format': 'mp4'
        }

        with yt_dlp.YoutubeDL(opcoes) as ydl:
            info = ydl.extract_info(url, download=True)
            titulo = info.get('title', 'Vídeo')

        arquivo_baixado = destino / f"{titulo}.mp4"

        arquivo_corrigido = converter_audio_para_aac(arquivo_baixado)

        arquivo_baixado.unlink()
        arquivo_corrigido.rename(arquivo_baixado)

        msg = f"Vídeo '{titulo}' baixado e convertido com áudio AAC com sucesso em {destino}."

        if modo == 'voz':
            falar(msg)
        return msg

    except Exception as e:
        erro = f"Erro ao baixar vídeo: {str(e)}"
        if modo == 'voz':
            falar(erro)
        return erro

# ========== Baixar Audio ==========
def limpar_nome_arquivo(nome):
    nome = unicodedata.normalize('NFKD', nome).encode('ASCII', 'ignore').decode()
    nome = re.sub(r'[^\w.-]', '_', nome)
    return nome

def converter_para_mp3(caminho_arquivo: Path):
    pasta = caminho_arquivo.parent
    nome_limpo = limpar_nome_arquivo(caminho_arquivo.name)
    caminho_limpo = pasta / nome_limpo

    if caminho_arquivo != caminho_limpo:
        caminho_arquivo.rename(caminho_limpo)
        caminho_arquivo = caminho_limpo

    caminho_saida = caminho_arquivo.with_suffix('.mp3')
    cmd = [
        'ffmpeg',
        '-i', str(caminho_arquivo),
        '-vn',
        '-ar', '44100',
        '-ac', '2',
        '-b:a', '192k',
        '-y',
        str(caminho_saida)
    ]
    processo = subprocess.run(cmd, capture_output=True, text=True, encoding='utf-8', errors='replace')
    if processo.returncode != 0:
        raise RuntimeError(f'Erro ao converter para MP3: {processo.stderr}')
    return caminho_saida

def baixar_audio_youtube(url, username, modo='texto'):
    try:
        destino = Path.home() / "Documents" / "Áudios Download"
        destino.mkdir(parents=True, exist_ok=True)

        opcoes = {
            'format': 'bestaudio/best',
            'outtmpl': str(destino / '%(title)s.%(ext)s'),
            'postprocessors': []
        }

        with yt_dlp.YoutubeDL(opcoes) as ydl:
            info = ydl.extract_info(url, download=True)
            titulo = info.get('title', 'Áudio')

        ext = info.get('ext', 'webm')  
        arquivo_baixado = destino / f"{titulo}.{ext}"

        arquivo_convertido = converter_para_mp3(arquivo_baixado)

        arquivo_baixado.unlink()

        msg = f"Áudio '{titulo}' baixado e convertido para MP3 com sucesso em {destino}."

        if modo == 'voz':
            falar(msg)
        return msg

    except Exception as e:
        erro = f"Erro ao baixar áudio: {str(e)}"
        if modo == 'voz':
            falar(erro)
        return erro

# ========== Gravação de Tela ==========
def iniciar_gravacao_sistema(username=None):
    try:
        pyautogui.hotkey('winleft', 'shift', 'r')
        time.sleep(1)

        largura, altura = pyautogui.size()

        x_inicial, y_inicial = 10, 10
        x_final, y_final = largura - 10, altura - 10

        pyautogui.moveTo(x_inicial, y_inicial, duration=0.5)
        pyautogui.mouseDown()
        pyautogui.moveTo(x_final, y_final, duration=1)
        pyautogui.mouseUp()
        time.sleep(0.5)
        
        time.sleep(1)
        pyautogui.moveTo(879, 44, duration=0.5)
        pyautogui.mouseDown()
        time.sleep(0.1)
        pyautogui.mouseUp()

        return "Gravação iniciada."
    except Exception as e:
        return f"Erro ao iniciar gravação: {str(e)}"

def parar_gravacao_sistema(username=None):
    try:
        pyautogui.click(879,44, duration=0.5)
        time.sleep(1)
        return "Gravação parada."
    except Exception as e:
        return f"Erro ao parar gravação: {str(e)}"

# ========== Funções de imagens ==========
class ImageAnalyser:
    def __init__(self):
        self.client = None
        self._init_client()

    def _init_client(self):
        try:
            from ai_service import client
            self.client = client
        except Exception:
            self.client = None

    def analisar(
        self,
        image_path: str,
        session_id: str,
        username: str | None = None
    ) -> str:

        if not os.path.exists(image_path):
            return f"❌ Caminho inválido: {image_path}"

        if self.client is None:
            if not recarregar_llm():
                return "❌ Gemini indisponível."

            from ai_service import client
            self.client = client

        try:
            image = Image.open(image_path).convert("RGB")

            prompt_usuario = (
                "Analise a imagem a seguir considerando o contexto da conversa. "
                "Descreva objetivamente tudo o que for relevante."
            )

            # Histórico DA SESSÃO (mesmo padrão do chat)
            mensagens = construir_historico(
                session_id=session_id,
                input_usuario=prompt_usuario
            )

            # Injeta a imagem na última mensagem do usuário
            mensagens[-1] = {
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt_usuario},
                    {"type": "image", "image": image}
                ]
            }

            response = self.client.models.generate_content(
                model=MODEL_NAME,
                contents=mensagens,
                generation_config={
                    "temperature": 0.4
                }
            )

            resposta_texto = response.text.strip()

            # Persistência POR SESSÃO
            adicionar_mensagem_chat(
                session_id,
                f"[IMAGEM] {image_path}",
                "human"
            )
            adicionar_mensagem_chat(
                session_id,
                resposta_texto,
                "ai"
            )

            # Auditoria opcional
            if username:
                registrar_log(username, f"[{session_id}] Análise de imagem: {image_path}")
                registrar_log(username, f"[{session_id}] Resultado: {resposta_texto}")

            return resposta_texto

        except Exception as e:
            erro = str(e)
            if username:
                registrar_log(username, f"[{session_id}] Erro imagem: {erro}")
            return f"❌ Erro ao analisar imagem: {erro}"

# =====================================================
# COMANDO
# =====================================================

def analisar_imagem_comando(
    caminho: str,
    session_id: str,
    username: str | None = None,
    modo: str = "texto"
) -> str:

    if not os.path.exists(caminho):
        return f"❌ Caminho inválido: {caminho}"

    analyser = ImageAnalyser()
    resultado = analyser.analisar(
        image_path=caminho,
        session_id=session_id,
        username=username
    )

    if modo == "voz":
        falar(resultado)

    return resultado
# ========== Funções de agenda ==========

AGENDA_DIR = Path.home() / "Documents" / "Agenda"
AGENDA_DIR.mkdir(parents=True, exist_ok=True)
COLUNAS = ["Tarefa", "DataHora", "Status"]

def _sanitize_username(username: str) -> str:
    return re.sub(r"[^a-zA-Z0-9_-]", "", username.lower())

def get_agenda_path(username: str) -> Path:
    return AGENDA_DIR / f"agenda_{_sanitize_username(username)}.xlsx"

def inicializar_agenda(username: str) -> None:
    path = get_agenda_path(username)
    if not path.exists():
        df = pd.DataFrame(columns=COLUNAS)
        df.to_excel(path, index=False)

def ler_agenda_df(username: str) -> pd.DataFrame:
    inicializar_agenda(username)
    path = get_agenda_path(username)
    
    try:
        df = pd.read_excel(path)
        
        # Garantir que as colunas existam
        for coluna in COLUNAS:
            if coluna not in df.columns:
                df[coluna] = None
        
        if not df.empty and "DataHora" in df.columns:
            df["DataHora"] = pd.to_datetime(df["DataHora"], errors="coerce")
        
        return df
    except Exception as e:
        print(f"Erro ao ler agenda: {e}")
        # Retornar dataframe vazio em caso de erro
        return pd.DataFrame(columns=COLUNAS)

def salvar_agenda_df(df: pd.DataFrame, username: str) -> None:
    try:
        df.to_excel(get_agenda_path(username), index=False)
    except Exception as e:
        print(f"Erro ao salvar agenda: {e}")

def _parse_datetime(data: str, hora: str | None) -> datetime:
    if hora:
        return datetime.strptime(f"{data} {hora}", "%d/%m/%Y %H:%M")
    return datetime.strptime(data, "%d/%m/%Y")

# =========================
# FUNÇÕES CRUD DA AGENDA
# =========================
def adicionar_tarefa(tarefa: str, data: str, hora: str | None, username: str) -> str:
    """Função base para adicionar tarefa (usada pelas outras funções)"""
    try:
        dt = _parse_datetime(data, hora)
    except ValueError as e:
        return f"❌ Data ou hora inválida. Use DD/MM/AAAA e HH:MM. Erro: {e}"

    df = ler_agenda_df(username)

    nova = pd.DataFrame([{
        "Tarefa": tarefa.strip(),
        "DataHora": dt,
        "Status": "Pendente"
    }])

    df = pd.concat([df, nova], ignore_index=True)
    salvar_agenda_df(df, username)

    hora_str = f" às {dt.strftime('%H:%M')}" if hora else ""
    return f"✅ Tarefa '{tarefa}' adicionada para {dt.strftime('%d/%m/%Y')}{hora_str}"

def adicionar_tarefa_interativa(match, username, modo='texto'):
    """Versão interativa que pede os dados via input"""
    try:
        print(f"\n{Colors.CYAN}📝 Adicionar nova tarefa{Colors.RESET}")
        print(f"{Colors.GRAY}{'='*40}{Colors.RESET}")
        
        tarefa = input(f"{Colors.PURPLE}>{Colors.RESET} Descrição da tarefa: ").strip()
        if not tarefa:
            return "❌ Tarefa não pode ser vazia."
        
        data = input(f"{Colors.PURPLE}>{Colors.RESET} Data (DD/MM/AAAA): ").strip()
        if not data:
            return "❌ Data é obrigatória."
        
        hora = input(f"{Colors.PURPLE}>{Colors.RESET} Hora (HH:MM ou Enter para sem hora): ").strip()
        if not hora:
            hora = None
        
        resultado = adicionar_tarefa(tarefa, data, hora, username)
        
        if modo == 'voz':
            if "✅" in resultado:
                falar("Tarefa adicionada com sucesso")
            else:
                falar("Erro ao adicionar tarefa")
        
        return resultado
        
    except Exception as e:
        return f"❌ Erro ao adicionar tarefa: {e}"

def listar_agenda(username, modo='texto'):
    """Lista todas as tarefas da agenda"""
    df = ler_agenda_df(username)

    if df.empty:
        mensagem = "📭 Agenda vazia."
        if modo == 'voz':
            falar("Sua agenda está vazia")
        return mensagem

    # Ordenar por data
    df = df.sort_values('DataHora')
    
    linhas = [f"\n{Colors.CYAN}📅 Agenda de {username}{Colors.RESET}"]
    linhas.append(f"{Colors.GRAY}{'='*40}{Colors.RESET}")

    agora = datetime.now()
    atrasadas = 0
    concluidas = 0
    
    for i, row in df.iterrows():
        # Converter para datetime se for string
        if isinstance(row["DataHora"], str):
            try:
                data_hora = datetime.strptime(row["DataHora"], "%Y-%m-%d %H:%M:%S")
            except:
                data_hora = agora
        else:
            data_hora = row["DataHora"]
        
        status = "✅" if row["Status"] == "Concluído" else "⏳"
        
        # Verificar se está atrasada
        if row["Status"] != "Concluído" and data_hora < agora:
            status = f"{Colors.RED}⚠{Colors.RESET}"
            atrasadas += 1
        
        if row["Status"] == "Concluído":
            concluidas += 1
            
        data_str = data_hora.strftime("%d/%m/%Y %H:%M") if pd.notna(data_hora) else "Sem data"
        linhas.append(f"{i+1}. {status} {row['Tarefa']} — {data_str}")

    # Adicionar resumo
    total = len(df)
    linhas.append(f"\n{Colors.GRAY}{'='*40}{Colors.RESET}")
    linhas.append(f"📊 Resumo: {total} tarefas")
    linhas.append(f"   ✅ Concluídas: {concluidas}")
    linhas.append(f"   ⏳ Pendentes: {total - concluidas}")
    
    if atrasadas > 0:
        linhas.append(f"   {Colors.RED}⚠ Atrasadas: {atrasadas}{Colors.RESET}")
    
    resultado = "\n".join(linhas)
    
    if modo == 'voz':
        falar(f"Sua agenda tem {total} tarefas. {concluidas} concluídas, {atrasadas} atrasadas")
    
    return resultado

def marcar_como_concluida(termo: str, username: str) -> str:
    """Função base para marcar tarefa como concluída"""
    df = ler_agenda_df(username)

    if df.empty:
        return "❌ Agenda vazia."

    mask = df["Tarefa"].str.lower().str.contains(termo.lower(), na=False)
    encontrados = df[mask]

    if encontrados.empty:
        return f"❌ Nenhuma tarefa encontrada com '{termo}'."

    if len(encontrados) > 1:
        # Mostrar múltiplas correspondências
        linhas = ["⚠ Mais de uma tarefa encontrada. Selecione:"]
        for i, (idx, row) in enumerate(encontrados.iterrows(), 1):
            data = row["DataHora"].strftime("%d/%m/%Y %H:%M") if pd.notna(row["DataHora"]) else "Sem data"
            linhas.append(f"{i}. {row['Tarefa']} — {data}")
        return "\n".join(linhas)

    idx = encontrados.index[0]
    df.at[idx, "Status"] = "Concluído"
    salvar_agenda_df(df, username)
    return f"✅ Tarefa '{df.at[idx, 'Tarefa']}' marcada como concluída."

def marcar_como_concluida_comando(match, username, modo='texto'):
    """Wrapper para extrair termo do comando ou modo interativo"""
    termo = ""
    
    if match.lastindex >= 1:
        termo = match.group(1).strip()
    
    if not termo:
        # Modo interativo
        print(f"\n{Colors.CYAN}✅ Marcar tarefa como concluída{Colors.RESET}")
        df = ler_agenda_df(username)
        
        if df.empty:
            return "📭 Agenda vazia."
        
        # Mostrar apenas tarefas pendentes
        pendentes = df[df["Status"] != "Concluído"]
        if pendentes.empty:
            return "🎉 Todas as tarefas já estão concluídas!"
        
        print(f"\n{Colors.YELLOW}Tarefas pendentes:{Colors.RESET}")
        for i, row in pendentes.iterrows():
            data = row["DataHora"].strftime("%d/%m/%Y %H:%M") if pd.notna(row["DataHora"]) else "Sem data"
            print(f"{i+1}. {row['Tarefa']} — {data}")
        
        try:
            escolha = input(f"\n{Colors.PURPLE}>{Colors.RESET} Número da tarefa ou nome: ").strip()
            
            # Tentar como número
            if escolha.isdigit():
                num = int(escolha)
                if 1 <= num <= len(pendentes):
                    tarefa = pendentes.iloc[num-1]["Tarefa"]
                    termo = tarefa
                else:
                    return "❌ Número inválido."
            else:
                # Usar como termo de busca
                termo = escolha
        except:
            termo = input(f"{Colors.PURPLE}>{Colors.RESET} Nome da tarefa: ").strip()
    
    resultado = marcar_como_concluida(termo, username)
    
    if modo == 'voz':
        if "✅" in resultado:
            falar("Tarefa marcada como concluída")
        elif "⚠" in resultado:
            falar("Múltiplas tarefas encontradas")
        else:
            falar("Erro ao marcar tarefa")
    
    return resultado

def remover_tarefa(termo: str, username: str) -> str:
    """Função base para remover tarefa"""
    df = ler_agenda_df(username)

    if df.empty:
        return "❌ Agenda vazia."

    mask = df["Tarefa"].str.lower().str.contains(termo.lower(), na=False)
    encontrados = df[mask]

    if encontrados.empty:
        return f"❌ Nenhuma tarefa encontrada com '{termo}'."

    if len(encontrados) > 1:
        # Mostrar múltiplas correspondências
        linhas = ["⚠ Mais de uma tarefa encontrada. Selecione:"]
        for i, (idx, row) in enumerate(encontrados.iterrows(), 1):
            status = "✅" if row["Status"] == "Concluído" else "⏳"
            data = row["DataHora"].strftime("%d/%m/%Y %H:%M") if pd.notna(row["DataHora"]) else "Sem data"
            linhas.append(f"{i}. {status} {row['Tarefa']} — {data}")
        return "\n".join(linhas)

    idx = encontrados.index[0]
    tarefa = df.at[idx, "Tarefa"]

    df = df.drop(index=idx)
    salvar_agenda_df(df, username)

    return f"🗑 Tarefa '{tarefa}' removida."

def remover_tarefa_comando(match, username, modo='texto'):
    """Wrapper para remover tarefa (interativo ou por comando)"""
    termo = ""
    
    if match.lastindex >= 1:
        termo = match.group(1).strip()
    
    if not termo:
        # Modo interativo
        print(f"\n{Colors.CYAN}🗑 Remover tarefa{Colors.RESET}")
        df = ler_agenda_df(username)
        
        if df.empty:
            return "📭 Agenda vazia."
        
        print(f"\n{Colors.YELLOW}Todas as tarefas:{Colors.RESET}")
        for i, row in df.iterrows():
            status = "✅" if row["Status"] == "Concluído" else "⏳"
            data = row["DataHora"].strftime("%d/%m/%Y %H:%M") if pd.notna(row["DataHora"]) else "Sem data"
            print(f"{i+1}. {status} {row['Tarefa']} — {data}")
        
        try:
            escolha = input(f"\n{Colors.PURPLE}>{Colors.RESET} Número da tarefa para remover: ").strip()
            
            # Tentar como número
            if escolha.isdigit():
                num = int(escolha)
                if 1 <= num <= len(df):
                    tarefa = df.iloc[num-1]["Tarefa"]
                    termo = tarefa
                else:
                    return "❌ Número inválido."
            else:
                # Usar como termo de busca
                termo = escolha
        except:
            termo = input(f"{Colors.PURPLE}>{Colors.RESET} Nome da tarefa para remover: ").strip()
    
    resultado = remover_tarefa(termo, username)
    
    if modo == 'voz':
        if "🗑" in resultado:
            falar("Tarefa removida")
        elif "⚠" in resultado:
            falar("Múltiplas tarefas encontradas")
        else:
            falar("Erro ao remover tarefa")
    
    return resultado

def limpar_agenda_completa(username, modo='texto'):
    """Remove todas as tarefas da agenda"""
    confirmar = input(f"{Colors.RED}⚠ Tem certeza que deseja limpar TODA a agenda? (s/n): {Colors.RESET}").strip().lower()
    if confirmar in ['s', 'sim', 'y', 'yes']:
        try:
            path = get_agenda_path(username)
            if path.exists():
                path.unlink()
            inicializar_agenda(username)
            
            resultado = "🗑 Agenda completamente limpa!"
            if modo == 'voz':
                falar("Agenda limpa com sucesso")
            return resultado
        except Exception as e:
            return f"❌ Erro ao limpar agenda: {e}"
    else:
        return "❌ Operação cancelada."

def agenda_hoje(username, modo='texto'):
    """Mostra tarefas para hoje COM NOTIFICAÇÃO"""
    hoje = datetime.now().date()
    
    df = ler_agenda_df(username)
    if df.empty:
        mensagem = "📭 Agenda vazia."
        if modo == 'voz':
            falar("Sua agenda está vazia")
        return mensagem
    
    # Filtrar tarefas de hoje
    df['Data'] = df['DataHora'].dt.date
    hoje_tarefas = df[df['Data'] == hoje]
    
    if hoje_tarefas.empty:
        mensagem = "🎉 Nenhuma tarefa para hoje!"
        if modo == 'voz':
            falar("Você não tem tarefas para hoje")
        
        # Notificação positiva
        try:
            notification.notify(
                title='🎉 JARVIS - Sem Tarefas!',
                message='Você não tem tarefas para hoje!',
                app_name='JARVIS Assistant',
                timeout=5
            )
        except:
            pass
            
        return mensagem
    
    # Separar concluídas e pendentes
    concluidas = hoje_tarefas[hoje_tarefas["Status"] == "Concluído"]
    pendentes = hoje_tarefas[hoje_tarefas["Status"] != "Concluído"]
    
    linhas = [f"\n{Colors.CYAN}📅 Tarefas para hoje ({hoje.strftime('%d/%m/%Y')}){Colors.RESET}"]
    linhas.append(f"{Colors.GRAY}{'='*40}{Colors.RESET}")
    
    if not pendentes.empty:
        linhas.append(f"{Colors.YELLOW}⏳ Pendentes:{Colors.RESET}")
        for i, row in pendentes.iterrows():
            hora = row["DataHora"].strftime("%H:%M") if pd.notna(row["DataHora"]) else "Dia todo"
            linhas.append(f"  • {row['Tarefa']} — {hora}")
        
        # Notificação com tarefas pendentes
        try:
            tasks_list = "\n".join([f"• {row['Tarefa']}" for _, row in pendentes.iterrows()[:3]])
            if len(pendentes) > 3:
                tasks_list += f"\n• e mais {len(pendentes) - 3} tarefas..."
            
            notification.notify(
                title=f'📋 JARVIS - {len(pendentes)} Tarefa(s) Pendente(s)',
                message=tasks_list,
                app_name='JARVIS Assistant',
                timeout=15
            )
        except:
            pass
    
    if not concluidas.empty:
        linhas.append(f"\n{Colors.GREEN}✅ Concluídas:{Colors.RESET}")
        for i, row in concluidas.iterrows():
            hora = row["DataHora"].strftime("%H:%M") if pd.notna(row["DataHora"]) else "Dia todo"
            linhas.append(f"  • {row['Tarefa']} — {hora}")
    
    # Resumo
    linhas.append(f"\n{Colors.GRAY}{'='*40}{Colors.RESET}")
    linhas.append(f"📊 Hoje: {len(pendentes)} pendentes, {len(concluidas)} concluídas")
    
    resultado = "\n".join(linhas)
    
    if modo == 'voz':
        if len(pendentes) > 0:
            falar(f"Você tem {len(pendentes)} tarefas pendentes para hoje")
        else:
            falar("Todas as tarefas de hoje estão concluídas")
    
    return resultado

def editar_tarefa(match, username, modo='texto'):
    """Edita uma tarefa existente"""
    try:
        print(f"\n{Colors.CYAN}✏️ Editar tarefa{Colors.RESET}")
        df = ler_agenda_df(username)
        
        if df.empty:
            return "📭 Agenda vazia."
        
        print(f"\n{Colors.YELLOW}Todas as tarefas:{Colors.RESET}")
        for i, row in df.iterrows():
            status = "✅" if row["Status"] == "Concluído" else "⏳"
            data = row["DataHora"].strftime("%d/%m/%Y %H:%M") if pd.notna(row["DataHora"]) else "Sem data"
            print(f"{i+1}. {status} {row['Tarefa']} — {data}")
        
        escolha = input(f"\n{Colors.PURPLE}>{Colors.RESET} Número da tarefa para editar: ").strip()
        
        if not escolha.isdigit():
            return "❌ Por favor, digite um número."
        
        num = int(escolha)
        if not (1 <= num <= len(df)):
            return "❌ Número inválido."
        
        idx = num - 1
        
        print(f"\n{Colors.GRAY}Deixe em branco para manter o valor atual{Colors.RESET}")
        
        # Editar descrição
        atual_desc = df.at[idx, "Tarefa"]
        nova_desc = input(f"Nova descrição (atual: '{atual_desc}'): ").strip()
        if nova_desc:
            df.at[idx, "Tarefa"] = nova_desc
        
        # Editar data/hora
        atual_data = df.at[idx, "DataHora"]
        if pd.notna(atual_data):
            data_str = atual_data.strftime("%d/%m/%Y")
            hora_str = atual_data.strftime("%H:%M")
        else:
            data_str = "Sem data"
            hora_str = ""
        
        nova_data = input(f"Nova data DD/MM/AAAA (atual: {data_str}): ").strip()
        nova_hora = input(f"Nova hora HH:MM (atual: {hora_str}): ").strip()
        
        if nova_data:
            try:
                if nova_hora:
                    nova_dt = _parse_datetime(nova_data, nova_hora)
                else:
                    nova_dt = _parse_datetime(nova_data, None)
                df.at[idx, "DataHora"] = nova_dt
            except ValueError as e:
                return f"❌ Data/hora inválida: {e}"
        
        # Editar status
        atual_status = df.at[idx, "Status"]
        novo_status = input(f"Status (C)oncluído/(P)endente (atual: {atual_status}): ").strip().lower()
        if novo_status in ['c', 'concluido', 'concluído']:
            df.at[idx, "Status"] = "Concluído"
        elif novo_status in ['p', 'pendente']:
            df.at[idx, "Status"] = "Pendente"
        
        salvar_agenda_df(df, username)
        
        resultado = f"✅ Tarefa '{df.at[idx, 'Tarefa']}' editada com sucesso."
        
        if modo == 'voz':
            falar("Tarefa editada com sucesso")
        
        return resultado
        
    except Exception as e:
        return f"❌ Erro ao editar tarefa: {e}"

def notificar_tarefas_do_dia(username):
    """Envia notificação com resumo das tarefas do dia"""
    try:
        df = ler_agenda_df(username)
        if df.empty:
            return
        
        hoje = datetime.now().date()
        
        # Filtrar tarefas de hoje
        df['Data'] = df['DataHora'].dt.date
        hoje_tarefas = df[df['Data'] == hoje]
        
        if not hoje_tarefas.empty:
            concluidas = len(hoje_tarefas[hoje_tarefas["Status"] == "Concluído"])
            pendentes = len(hoje_tarefas) - concluidas
            
            if pendentes > 0:
                notification.notify(
                    title='📅 JARVIS - Tarefas de Hoje',
                    message=f'{pendentes} tarefa(s) pendente(s) para hoje',
                    app_name='JARVIS Assistant',
                    timeout=10
                )
    except Exception as e:
        print(f"Erro ao notificar tarefas: {e}")

def checar_tarefas_atrasadas(username, modo='texto'):
    """Verifica e gerencia tarefas atrasadas COM NOTIFICAÇÕES"""
    df = ler_agenda_df(username)

    if df.empty:
        if modo == 'voz':
            falar("Sua agenda está vazia")
        return "📭 Agenda vazia."

    agora = datetime.now()

    atrasadas = df[
        (df["Status"] != "Concluído") &
        (df["DataHora"] < agora)
    ]

    if atrasadas.empty:
        if modo == 'voz':
            falar("Você não tem tarefas atrasadas")
        return "✅ Nenhuma tarefa atrasada!"

    # 📢 EXIBIR NOTIFICAÇÃO DO WINDOWS
    try:
        notification.notify(
            title='🚨 JARVIS - Tarefas Atrasadas!',
            message=f'Você tem {len(atrasadas)} tarefa(s) atrasada(s)',
            app_name='JARVIS Assistant',
            timeout=10,  # 10 segundos
            toast=False  # Para notificação padrão do Windows
        )
    except Exception as e:
        print(f"⚠ Não foi possível exibir notificação: {e}")

    mensagem = f"\n{Colors.RED}⚠ {len(atrasadas)} TAREFA(S) ATRASADA(S) ⚠{Colors.RESET}"
    
    if modo == 'voz':
        falar(f"ATENÇÃO! Você tem {len(atrasadas)} tarefas atrasadas!")
    
    print(mensagem)
    
    resultados = []
    
    for idx, row in atrasadas.iterrows():
        print(f"\n{Colors.GRAY}{'='*40}{Colors.RESET}")
        print(f"{Colors.RED}Tarefa atrasada:{Colors.RESET} {row['Tarefa']}")
        print(f"{Colors.RED}Devia ser feita em:{Colors.RESET} {row['DataHora'].strftime('%d/%m/%Y %H:%M')}")
        print(f"\n{Colors.YELLOW}O que deseja fazer?{Colors.RESET}")
        print("1. Marcar como concluída")
        print("2. Remarcar para nova data")
        print("3. Remover tarefa")
        print("4. Manter como está")

        escolha = input(f"{Colors.PURPLE}> Escolha (1-4): {Colors.RESET}").strip()

        if escolha == "1":
            df.at[idx, "Status"] = "Concluído"
            resultados.append(f"✅ '{row['Tarefa']}' marcada como concluída")
            if modo == 'voz':
                falar(f"Tarefa {row['Tarefa']} marcada como concluída")

        elif escolha == "2":
            nova_data = input("Nova data (DD/MM/AAAA): ").strip()
            nova_hora = input("Nova hora (HH:MM ou Enter): ").strip() or None

            try:
                nova_dt = _parse_datetime(nova_data, nova_hora)
                df.at[idx, "DataHora"] = nova_dt
                resultados.append(f"📅 '{row['Tarefa']}' remarcada para {nova_dt.strftime('%d/%m/%Y %H:%M')}")
                if modo == 'voz':
                    falar(f"Tarefa {row['Tarefa']} remarcada")
            except ValueError:
                resultados.append(f"❌ Data inválida para '{row['Tarefa']}'. Mantida como está.")

        elif escolha == "3":
            confirmar = input(f"Tem certeza que deseja remover '{row['Tarefa']}'? (s/n): ").strip().lower()
            if confirmar in ['s', 'sim', 'y', 'yes']:
                df = df.drop(index=idx)
                resultados.append(f"🗑 '{row['Tarefa']}' removida")
                if modo == 'voz':
                    falar(f"Tarefa {row['Tarefa']} removida")
            else:
                resultados.append(f"❌ '{row['Tarefa']}' não removida")

    if resultados:
        salvar_agenda_df(df, username)
        resultados.insert(0, f"\n{Colors.GREEN}✅ Ações realizadas:{Colors.RESET}")
        return "\n".join(resultados)
    
    return "⚠ Nenhuma ação realizada nas tarefas atrasadas."

def inicializar_sistema_agenda(username):
    """Inicializa o sistema de agenda para o usuário COM NOTIFICAÇÕES PERIÓDICAS"""
    # Criar diretório se não existir
    AGENDA_DIR.mkdir(parents=True, exist_ok=True)
    
    # Inicializar agenda
    inicializar_agenda(username)
    
    # Verificar tarefas atrasadas no início
    print(f"{Colors.YELLOW}⏳ Verificando tarefas atrasadas...{Colors.RESET}")
    
    # Criar thread para verificar tarefas periodicamente
    def verificar_periodicamente():
        while True:
            try:
                # Verificar a cada 30 minutos
                time.sleep(1800)  # 1800 segundos = 30 minutos
                
                # Verificar tarefas atrasadas
                df = ler_agenda_df(username)
                if df.empty:
                    continue
                
                agora = datetime.now()
                atrasadas = df[
                    (df["Status"] != "Concluído") &
                    (df["DataHora"] < agora)
                ]
                
                if not atrasadas.empty:
                    # Exibir notificação
                    notification.notify(
                        title='⏰ JARVIS - Lembrete de Tarefas',
                        message=f'Você tem {len(atrasadas)} tarefa(s) atrasada(s)',
                        app_name='JARVIS Assistant',
                        timeout=10
                    )
                    
                    # Falar alerta se modo voz estiver ativo
                    if modo == 'voz':
                        falar(f"Lembrete: você ainda tem {len(atrasadas)} tarefas atrasadas")
                        
            except Exception as e:
                print(f"Erro na verificação periódica: {e}")
                continue
    
    # Iniciar thread em segundo plano
    threading.Thread(target=verificar_periodicamente, daemon=True).start()
    
    # Verificar tarefas atrasadas agora
    resultado = checar_tarefas_atrasadas(username, modo='texto')
    
    # Mostrar tarefas de hoje
    hoje = agenda_hoje(username, modo='texto')
    if "Nenhuma tarefa para hoje" not in hoje:
        print(hoje)
    
    return resultado

# ========== Abrir sites ==========
def abrir_site(match, username):
    comando = match.group(0).lower()
    sites = {
        "github": "https://github.com/",
        "netflix": "https://www.netflix.com",
        "youtube": "https://youtube.com",
        "microsoft teams": "https://teams.microsoft.com",
        "teams": "https://teams.microsoft.com",
        "instagram": "https://www.instagram.com",
        "whatsapp": "https://web.whatsapp.com",
        "tik tok": "https://www.tiktok.com",
        "tiktok": "https://www.tiktok.com",
        "e-mail": "https://mail.google.com",
        "email": "https://mail.google.com",
        "calendario": "https://calendar.google.com/calendar/u/0/r",
        "meet": "https://meet.google.com/landing?hs=197&authuser=0",
        "google meet": "https://meet.google.com/landing?hs=197&authuser=0",
        "drive": "https://drive.google.com/drive/u/0/home",
        "google drive": "https://drive.google.com/drive/u/0/home"
    }
    for nome, url in sites.items():
        if nome in comando:
            try:
                webbrowser.open(url)
                return f"Abrindo {nome}, senhor."
            except Exception as e:
                return f"Erro ao abrir site {nome}: {e}"
    return "Site não reconhecido, senhor."

# ======= Atualizações e limpeza =======

def verificar_atualizacoes(match, username):
    try:
        subprocess.run("powershell -Command \"Get-WindowsUpdate\"", shell=True)
        return "Verificando atualizações do sistema, senhor."
    except Exception as e:
        return f"Erro ao verificar atualizações: {e}"

def atualizar_sistema(match, username):
    try:
        subprocess.run("powershell -Command \"Install-WindowsUpdate -AcceptAll -AutoReboot\"", shell=True)
        return "Atualizações sendo instaladas, senhor. O sistema pode reiniciar automaticamente."
    except Exception as e:
        return f"Erro ao atualizar sistema: {e}"

def limpar_lixo(match, username):
    try:
        pastas = [
            os.getenv('TEMP'),
            os.path.join(os.getenv('SystemRoot'), 'Temp')
        ]
        for pasta in pastas:
            for arquivo in os.listdir(pasta):
                caminho = os.path.join(pasta, arquivo)
                try:
                    if os.path.isfile(caminho) or os.path.islink(caminho):
                        os.unlink(caminho)
                    elif os.path.isdir(caminho):
                        shutil.rmtree(caminho)
                except Exception:
                    pass
        return "Arquivos temporários e lixo digital limpos com sucesso, senhor."
    except Exception as e:
        return f"Erro ao limpar arquivos: {e}"

# ========== Funções de data e hora ==========
def falar_hora(match, username):
    hora = datetime.now().strftime('%H:%M')
    if modo == 'voz':
        falar(f"Agora são {hora}")
    return f"Agora são {hora}"

def falar_data(match, username):
    data = datetime.now().strftime('%d/%m/%Y')
    if modo == 'voz':
        falar(f"Hoje é dia {data}")
    return f"Hoje é dia {data}"

# ========== Funções de pastas ==========
def encontrar_pasta(nome_pasta_usuario):
    home = os.path.expanduser("~")
    mapeamento = {
        "documentos": "Documents",
        "downloads": "Downloads",
        "imagens": "Pictures",
        "desktop": "Desktop",
        "musicas": "Music",
        "videos": "Videos"
    }
    nome_sistema = mapeamento.get(nome_pasta_usuario.lower(), nome_pasta_usuario)
    caminho = os.path.join(home, nome_sistema)
    if os.path.exists(caminho):
        return caminho
    for item in os.listdir(home):
        if nome_pasta_usuario.lower() in item.lower():
            caminho_tentativa = os.path.join(home, item)
            if os.path.isdir(caminho_tentativa):
                return caminho_tentativa
    return None

def abrir_pasta(match, username):
    nome_pasta_usuario = match.group(1).strip()
    caminho = encontrar_pasta(nome_pasta_usuario)
    if caminho:
        try:
            comando = f"Start-Process explorer -ArgumentList '{caminho}' -WindowStyle Maximized"
            subprocess.run(["powershell", "-Command", comando], check=True)
            return f"Abrindo a pasta {os.path.basename(caminho)}, senhor."
        except Exception as e:
            return f"Erro ao abrir a pasta: {e}"
    return f"Pasta '{nome_pasta_usuario}' não encontrada, senhor."

# ========== Listagem de sites ==========
def listar_sites(match, username):
    sites = {
        "github": "https://github.com/",
        "netflix": "https://www.netflix.com",
        "youtube": "https://youtube.com",
        "microsoft teams": "https://teams.microsoft.com",
        "instagram": "https://www.instagram.com",
        "whatsapp": "https://web.whatsapp.com",
        "tik tok": "https://www.tiktok.com",
        "e-mail": "https://mail.google.com",
        "calendario": "https://calendar.google.com/calendar/u/0/r",
        "meet": "https://meet.google.com/landing?hs=197&authuser=0",
        "google meet": "https://meet.google.com/landing?hs=197&authuser=0",
        "drive": "https://drive.google.com/drive/u/0/home",
        "google drive": "https://drive.google.com/drive/u/0/home"
    }
    return "Sites disponíveis:\n" + "\n".join(f"- {k}" for k in sites.keys())

# ========== Criação e manipulação de arquivos ==========
def criar_arquivo(match, username):
    documentos = Path.home() / "Documents"
    nome = input(f"{Colors.PURPLE}>{Colors.RESET} Digite o nome do arquivo (ex: texto.txt): ").strip()
    if not nome:
        falar("Nome de arquivo inválido.")
        return "Operação cancelada."
    conteudo = input(f"{Colors.PURPLE}>{Colors.RESET} Digite o conteúdo que deseja salvar: ")
    caminho_completo = documentos / nome
    try:
        with open(caminho_completo, "w", encoding="utf-8") as f:
            f.write(conteudo)
        return f"Arquivo '{nome}' criado com sucesso na pasta Documentos."
    except Exception as e:
        return f"Erro ao criar o arquivo: {e}"

def criar_codigo(match, username):
    documentos = Path.home() / "Documents"
    documentos.mkdir(exist_ok=True)
    linguagem = input(f"{Colors.PURPLE}>{Colors.RESET} Qual linguagem de programação você quer usar? ").strip().lower()
    descricao = input(f"{Colors.PURPLE}>{Colors.RESET} Descreva o que o código deve fazer: ").strip()
    prompt = f"Crie um código em {linguagem} que: {descricao}"
    try:
        codigo = gerar_resposta_ia(prompt, username)
    except Exception as e:
        return f"Erro ao gerar código com Gemini: {e}"
    extensoes = {
        "python": ".py",
        "javascript": ".js",
        "java": ".java",
        "html": ".html",
        "c": ".c",
        "cpp": ".cpp",
        "go": ".go",
        "php": ".php",
        "ruby": ".rb",
        "kotlin": ".kt",
        "swift": ".swift",
        "rust": ".rs",
        "csharp": ".cs",
        "c#": ".cs",
        "css": ".css",
        "sql": ".sql",
        "r": ".r"
    }
    ext = extensoes.get(linguagem, ".txt")
    nome_arquivo = input(f"{Colors.PURPLE}>{Colors.RESET} Nome do arquivo (sem extensão)? ").strip() + ext
    caminho = documentos / nome_arquivo
    try:
        with open(caminho, "w", encoding="utf-8") as f:
            f.write(codigo)
        return f"Código gerado e salvo em: {caminho}"
    except Exception as e:
        return f"Erro ao salvar o arquivo: {e}"

def listar_arquivos(match, username):
    extensao = (match.group(1) or "").lower().replace(".", "").strip()
    pasta = (match.group(2) or "Documentos").lower().strip()
    base_path = Path.home()
    nome_pasta = {
        "documentos": "Documents",
        "documents": "Documents",
        "area de trabalho": "Desktop",
        "área de trabalho": "Desktop",
        "desktop": "Desktop",
        "downloads": "Downloads"
    }.get(pasta, pasta)
    diretorio = base_path / nome_pasta
    if not diretorio.exists():
        return f"A pasta '{nome_pasta}' não foi encontrada, senhor."
    try:
        arquivos = list(diretorio.rglob(f"*.{extensao}")) if extensao else list(diretorio.rglob("*"))
        if not arquivos:
            if extensao:
                return f"Senhor, não encontrei arquivos '.{extensao}' na pasta '{nome_pasta}'."
            return f"Senhor, a pasta '{nome_pasta}' está vazia."
        lista = "\n- " + "\n- ".join([str(arq.relative_to(base_path)) for arq in arquivos])
        if extensao:
            return f"Senhor, encontrei os seguintes arquivos '.{extensao}' na pasta '{nome_pasta}' e suas subpastas:\n{lista}"
        return f"Senhor, encontrei os seguintes arquivos na pasta '{nome_pasta}' e suas subpastas:\n{lista}"
    except Exception as e:
        return f"Erro ao listar arquivos: {e}"

# ========== Funções para ler arquivos ==========
def buscar_arquivo_por_nome(nome_arquivo, pasta_base=None):
    if pasta_base is None:
        pasta_base = Path.home() / "Documents"
    for arquivo in pasta_base.rglob("*"):
        if arquivo.name.lower() == nome_arquivo.lower():
            return arquivo.resolve()
    return None

def ler_docx(caminho):
    try:
        doc = Document(caminho)
        return '\n'.join([p.text for p in doc.paragraphs])
    except Exception as e:
        return f"Erro ao ler .docx: {e}"

def ler_txt(caminho):
    try:
        with open(caminho, "r", encoding="utf-8") as f:
            return f.read()
    except Exception as e:
        return f"Erro ao ler TXT: {e}"

def ler_csv(caminho):
    try:
        df = pd.read_csv(caminho)
        return df.to_string(index=False)
    except Exception as e:
        return f"Erro ao ler CSV: {e}"

def ler_json(caminho):
    try:
        with open(caminho, "r", encoding="utf-8") as f:
            return json.load(f)
    except Exception as e:
        return f"Erro ao ler JSON: {e}"

def ler_pdf(caminho):
    try:
        doc = fitz.open(caminho)
        return ''.join(pagina.get_text() for pagina in doc)
    except Exception as e:
        return f"Erro ao ler PDF: {e}"

def ler_excel(caminho):
    try:
        df = pd.read_excel(caminho)
        return df.to_string(index=False)
    except Exception as e:
        return f"Erro ao ler Excel: {e}"

def ler_pptx(caminho):
    try:
        prs = Presentation(caminho)
        texto = ''
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texto += shape.text + "\n"
        return texto
    except Exception as e:
        return f"Erro ao ler PPTX: {e}"

def analisar_arquivos(match, username):
    try:
        nome_arquivo = match.group(2).strip()
        arquivo_encontrado = buscar_arquivo_por_nome(nome_arquivo)
        if not arquivo_encontrado:
            return f"Não encontrei o arquivo '{nome_arquivo}' na pasta Documentos, senhor."
        
        sufixo = arquivo_encontrado.suffix.lower()
        conteudo = None

        if sufixo == ".txt":
            conteudo = ler_txt(arquivo_encontrado)
        elif sufixo == ".docx":
            conteudo = ler_docx(arquivo_encontrado)
        elif sufixo == ".csv":
            conteudo = ler_csv(arquivo_encontrado)
        elif sufixo == ".json":
            conteudo = json.dumps(ler_json(arquivo_encontrado), indent=4)
        elif sufixo == ".pdf":
            conteudo = ler_pdf(arquivo_encontrado)
        elif sufixo in [".xlsx", ".xls"]:
            conteudo = ler_excel(arquivo_encontrado)
        elif sufixo == ".pptx":
            conteudo = ler_pptx(arquivo_encontrado)
        else:
            return f"Formato de arquivo '{sufixo}' não suportado, senhor."

        if not conteudo or not conteudo.strip():
            return "O arquivo está vazio ou ilegível, senhor."

        prompt = f"Analise esse conteúdo extraído do arquivo:\n\n{conteudo}"
        resposta = gerar_resposta_ia(prompt, username)
        return resposta
    except Exception as e:
        return f"Erro ao analisar arquivo: {e}"

# ========== Função fallback Gemini ==========
def responder_com_gemini_fallback(match, username):
    comando = match.group(0)
    return gerar_resposta_ia(comando, username)

# ========== Lista de comandos ==========
padroes = [
    # E-mail
    (re.compile(r'\/(?:enviar\s+)?e-?mail\b', re.IGNORECASE), 
     enviar_email),
    
    # Aplicativos
    (re.compile(r'\/listar\s+(?:apps|aplicativos)\b', re.IGNORECASE), 
     listar_aplicativos_winapps),
    
    (re.compile(r'\/info\s+(?:app|aplicativo)\s+(.+)', re.IGNORECASE), 
     lambda m, u: info_aplicativo_winapps(m.group(1).strip(), u)),
    
    (re.compile(r'\/desinstalar\s+(?:app|aplicativo)\s+(.+)', re.IGNORECASE), 
     lambda m, u: desinstalar_app_winapps(m.group(1).strip(), u)),

    # WhatsApp
    (re.compile(r'\/(?:enviar\s+)?whatsapp\s+grupo\b', re.IGNORECASE), 
     enviar_whatsapp_grupo),
    
    (re.compile(r'\/(?:enviar\s+)?whatsapp\s+agendado\b', re.IGNORECASE), 
    enviar_whatsapp_agendado),

    (re.compile(r'\/(?:enviar\s+)?whatsapp\b', re.IGNORECASE), 
    enviar_whatsapp),
    
    # YouTube e Pesquisa
    (re.compile(r'\/tocar\s+(?:no\s+)?youtube\b', re.IGNORECASE), 
     tocar_musica_pywhatkit),
    
    (re.compile(r'\/pesquisar\s+(.+?)(?:\s+no\s+google)?$', re.IGNORECASE), 
     pesquisar_google_pywhatkit),
    
    # Sites
    (re.compile(r'\/listar\s+sites\b', re.IGNORECASE), listar_sites),
    
    # Análises
    (re.compile(r'\/analisar\s+arquivo\s+(.+)', re.IGNORECASE), 
     lambda m, u: analisar_arquivos(m, u)),
    
    (re.compile(r'\/analisar\s+site\s+(.+)', re.IGNORECASE), 
     lambda m, u: analisar_site(m.group(1).strip(), u)),
    
    # Instalação/Desinstalação
    (re.compile(r"\/instalar\s+([a-zA-Z0-9\-\.]+)", re.IGNORECASE), 
     lambda m, u: instalar_programa_via_cmd_admin(m.group(1), u)),
    
    (re.compile(r"\/desinstalar\s+([a-zA-Z0-9\-\.]+)", re.IGNORECASE), 
     lambda m, u: desinstalar_programa(m.group(1), u, 'texto')),
    
    # Download YouTube
    (re.compile(r"\/baixar\s+video\s+(https?://[^\s]+)", re.IGNORECASE), 
     lambda m, u: baixar_video_youtube(m.group(1), u)),
    
    (re.compile(r"\/baixar\s+audio\s+(https?://[^\s]+)", re.IGNORECASE), 
     lambda m, u: baixar_audio_youtube(m.group(1), u)),
    
    # Gravação de tela
    (re.compile(r'\/gravar\s+tela\b', re.IGNORECASE), 
     lambda m, u: iniciar_gravacao_sistema()),
    
    (re.compile(r'\/parar\s+gravacao\b', re.IGNORECASE), 
     lambda m, u: parar_gravacao_sistema()),
    
    # Pastas
    (re.compile(r'\/abrir\s+pasta\s+(.+)', re.IGNORECASE), abrir_pasta),
    
    # Abrir sites
    (re.compile(r'\/abrir\s+(youtube|netflix|microsoft teams|github|instagram|tik\s*tok|tiktok|e-?mail|email|whatsapp|google met|calendario|meet|drive|google drive)\b', re.IGNORECASE), 
     abrir_site),
    
    # Abrir aplicativos
    (re.compile(r'\/abrir\s+(.+)', re.IGNORECASE), 
     abrir_aplicativo_winapps),
    
    # Análise de imagem
    (re.compile(r'\/analisar\s+imagem\s+(.+)', re.IGNORECASE), 
     lambda m, u: analisar_imagem_comando(m.group(1).strip(), u)),
    
     # Agenda
    (re.compile(r'\/(?:ler|ver)\s+agenda\b', re.IGNORECASE), 
    lambda m, u, modo='texto': listar_agenda(u, modo)),
    (re.compile(r'\/agenda\s+hoje\b', re.IGNORECASE), 
    lambda m, u, modo='texto': agenda_hoje(u, modo)),
    (re.compile(r'\/adicionar\s+tarefa$', re.IGNORECASE), 
    lambda m, u, modo='texto': adicionar_tarefa_interativa(m, u, modo)),
    (re.compile(r'\/marcar\s+(?:como\s+)?concluida\s+(.+)', re.IGNORECASE), 
    lambda m, u, modo='texto': marcar_como_concluida_comando(m, u, modo)),
    (re.compile(r'\/remover\s+tarefa\s+(.+)', re.IGNORECASE), 
    lambda m, u, modo='texto': remover_tarefa_comando(m, u, modo)),
    (re.compile(r'\/limpar\s+agenda\b', re.IGNORECASE), 
    lambda m, u, modo='texto': limpar_agenda_completa(u, modo)),
    (re.compile(r'\/editar\s+tarefa\b', re.IGNORECASE), 
    lambda m, u, modo='texto': editar_tarefa(m, u, modo)),
    (re.compile(r'\/tarefas\s+atrasadas\b', re.IGNORECASE), 
    lambda m, u, modo='texto': checar_tarefas_atrasadas(u, modo)),
    (re.compile(r'\/inicializar\s+agenda\b', re.IGNORECASE), 
    lambda m, u, modo='texto': inicializar_sistema_agenda(u)),
    
    # Sistema
    (re.compile(r'\/verificar\s+atualizacoes\b', re.IGNORECASE), verificar_atualizacoes),
    (re.compile(r'\/atualizar\s+sistema\b', re.IGNORECASE), atualizar_sistema),
    (re.compile(r'\/limpar\s+lixo\b', re.IGNORECASE), limpar_lixo),
    
    # Data e hora
    (re.compile(r'\/horas?\b', re.IGNORECASE), falar_hora),
    (re.compile(r'\/data\b', re.IGNORECASE), falar_data),
    
    # Arquivos
    (re.compile(r'\/listar\s+arquivos(?:\s+\.(\w+))?(?:\s+em\s+(.+))?', re.IGNORECASE), listar_arquivos),
    (re.compile(r'\/criar\s+(?:arquivo\s+)?texto\b', re.IGNORECASE), criar_arquivo),
    (re.compile(r'\/criar\s+codigo\b', re.IGNORECASE), criar_codigo),
]

# Variável global para modo
modo = 'texto'

# ========== Enhanced Command Processor ==========
def processar_comando(comando, username, token=None, modo="texto"):
    if token:
        session_id = obter_session_id_por_token(token)
        if session_id:
            return gerar_resposta_ia(comando, session_id, username)
    
    return "Comando não reconhecido. Como posso ajudar?"