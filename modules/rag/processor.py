import os
import re
import uuid
from pathlib import Path
from typing import Optional

from database.sqlite.connection import get_connection, release_connection

AREA_POR_ROLE = {
    "admin": "Administração",
    "tech": "Tecnologia",
    "security": "Segurança",
    "marketing": "Marketing",
    "finance": "Financeiro",
    "legal": "Jurídico",
    "rh": "Recursos Humanos",
    "user": "Geral",
}


def _departamento_do_usuario(username: str) -> str:
    from modules.permissions.rbac import get_user_roles

    roles = get_user_roles(username)
    for r in roles:
        area = AREA_POR_ROLE.get(r["name"])
        if area:
            return area
    return "Geral"


def _executar(query, params=None, *, fetch=False, fetchone=False, commit=False):
    conn = get_connection()
    try:
        cursor = conn.cursor()
        cursor.execute(query, params or ())
        if commit:
            conn.commit()
            return True
        if fetchone:
            row = cursor.fetchone()
            return dict(row) if row else None
        if fetch:
            return [dict(r) for r in cursor.fetchall()]
        return True
    finally:
        release_connection(conn)


CHUNK_SIZE = 500
CHUNK_OVERLAP = 50

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt", ".md"}


class DocumentProcessor:
    def extract_text(self, file_path: str) -> str:
        ext = Path(file_path).suffix.lower()
        if ext == ".pdf":
            return self._extract_pdf(file_path)
        elif ext == ".docx":
            return self._extract_docx(file_path)
        elif ext == ".pptx":
            return self._extract_pptx(file_path)
        elif ext == ".txt":
            return self._extract_txt(file_path)
        elif ext == ".md":
            return self._extract_txt(file_path)
        else:
            raise ValueError(f"Formato não suportado: {ext}")

    def _extract_pdf(self, file_path: str) -> str:
        try:
            import pymupdf

            doc = pymupdf.open(file_path)
            text = "\n".join(page.get_text() for page in doc)
            doc.close()
            return text
        except ImportError:
            try:
                from PyPDF2 import PdfReader

                reader = PdfReader(file_path)
                return "\n".join(page.extract_text() or "" for page in reader.pages)
            except ImportError:
                return ""

    def _extract_docx(self, file_path: str) -> str:
        try:
            from docx import Document

            doc = Document(file_path)
            return "\n".join(p.text for p in doc.paragraphs)
        except ImportError:
            return ""

    def _extract_pptx(self, file_path: str) -> str:
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            return "\n".join(texts)
        except ImportError:
            return ""

    def _extract_txt(self, file_path: str) -> str:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            return f.read()

    def chunk_text(self, text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
        text = re.sub(r"\s+", " ", text).strip()
        if not text:
            return []

        chunks = []
        start = 0
        while start < len(text):
            end = start + chunk_size
            if end >= len(text):
                chunks.append(text[start:].strip())
                break

            chunk = text[start:end]
            last_period = chunk.rfind(".")
            last_newline = chunk.rfind("\n")
            split_at = max(last_period + 1, last_newline + 1)

            if split_at <= start:
                split_at = end

            chunks.append(text[start:split_at].strip())
            start = split_at - overlap if split_at > overlap else split_at

        return [c for c in chunks if c]

    def process_and_store(
        self,
        file_path: str,
        username: str,
        filename: str,
        original_name: str,
        file_type: str,
        file_size: int,
        department: str = "",
    ) -> str:
        text = self.extract_text(file_path)
        chunks = self.chunk_text(text)

        if not department:
            department = _departamento_do_usuario(username)

        doc_id = str(uuid.uuid4())
        _executar(
            """
            INSERT INTO documents (id, filename, original_name, file_type, file_size, username, department, processed)
            VALUES (?, ?, ?, ?, ?, ?, ?, ?)
            """,
            (doc_id, filename, original_name, file_type, file_size, username, department, 1 if chunks else 0),
            commit=True,
        )

        for i, chunk_content in enumerate(chunks):
            chunk_id = str(uuid.uuid4())
            _executar(
                """
                INSERT INTO document_chunks (id, document_id, chunk_index, content)
                VALUES (?, ?, ?, ?)
                """,
                (chunk_id, doc_id, i, chunk_content),
                commit=True,
            )

        from modules.rag.engine import RAGEngine

        engine = RAGEngine()
        if chunks:
            engine.index_document(
                doc_id,
                chunks,
                metadata={
                    "filename": original_name,
                    "file_type": file_type,
                    "username": username,
                    "department": department,
                },
            )

        return doc_id

    def get_document(self, doc_id: str) -> Optional[dict]:
        return _executar(
            "SELECT id, filename, original_name, file_type, file_size, username, department, uploaded_at, processed FROM documents WHERE id = ?",
            (doc_id,),
            fetchone=True,
        )

    def list_documents(self, username: Optional[str] = None, department: Optional[str] = None) -> list[dict]:
        if department:
            return _executar(
                "SELECT id, filename, original_name, file_type, file_size, username, department, uploaded_at, processed FROM documents WHERE department = ? ORDER BY uploaded_at DESC",
                (department,),
                fetch=True,
            ) or []
        if username:
            return _executar(
                "SELECT id, filename, original_name, file_type, file_size, username, department, uploaded_at, processed FROM documents WHERE username = ? ORDER BY uploaded_at DESC",
                (username,),
                fetch=True,
            ) or []
        return _executar(
            "SELECT id, filename, original_name, file_type, file_size, username, department, uploaded_at, processed FROM documents ORDER BY uploaded_at DESC",
            fetch=True,
        ) or []

    def delete_document(self, doc_id: str):
        from modules.rag.engine import RAGEngine

        engine = RAGEngine()
        try:
            engine.delete_document(doc_id)
        except Exception:
            pass

        _executar(
            "DELETE FROM document_chunks WHERE document_id = ?",
            (doc_id,),
            commit=True,
        )
        _executar(
            "DELETE FROM documents WHERE id = ?",
            (doc_id,),
            commit=True,
        )
