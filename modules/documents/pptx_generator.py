import os
from pathlib import Path
from typing import Optional

from modules.documents.base import ensure_output_dir, safe_filename


def generate_pptx(
    title: str,
    slides: list[dict],
    filename: Optional[str] = None,
    author: str = "JARVIS",
    template_path: Optional[str] = None,
) -> str:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    if template_path and os.path.exists(template_path):
        prs = Presentation(template_path)
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

    def _add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=RGBColor(0x1A, 0x1A, 0x2E), alignment=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = alignment
        return tf

    def _add_bullet_frame(slide, left, top, width, height, items, font_size=14):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(font_size)
            p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            p.level = 0
            p.space_after = Pt(4)

    # Title slide
    title_slide = prs.slides.add_slide(blank_layout)
    _add_textbox(title_slide, 1, 2.5, 11.3, 2, title, font_size=36, bold=True, alignment=PP_ALIGN.CENTER)
    _add_textbox(title_slide, 1, 4.5, 11.3, 1, f"Gerado por JARVIS • {author}", font_size=14, color=RGBColor(0x66, 0x66, 0x66), alignment=PP_ALIGN.CENTER)

    # Content slides
    for slide_data in slides:
        slide_type = slide_data.get("type", "content")
        slide_title = slide_data.get("title", "")
        items = slide_data.get("items", [])
        content_text = slide_data.get("content", "")

        s = prs.slides.add_slide(blank_layout)

        _add_textbox(s, 0.5, 0.3, 12.3, 0.8, slide_title, font_size=28, bold=True)

        if items:
            _add_bullet_frame(s, 0.5, 1.3, 12.3, 5.5, items)
        elif content_text:
            _add_textbox(s, 0.5, 1.3, 12.3, 5.5, content_text, font_size=16)

    output_dir = ensure_output_dir("pptx")
    name = filename or safe_filename(title)
    if not name.endswith(".pptx"):
        name += ".pptx"

    output_path = output_dir / name
    prs.save(str(output_path))
    return str(output_path)
